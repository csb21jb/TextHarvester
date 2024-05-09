import os
import pptx
import pdfplumber
from docx import Document as DocxDocument
from tqdm import tqdm
import logging
from colorama import Fore, Style, init
import re
import textract

# Initialize colorama
init(autoreset=True)

# Set up logging
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')

# Function to sanitize text
def sanitize_text(text):
    return re.sub(r'[\x00-\x1F\x7F-\x9F]', ' ', text)

def extract_text_from_pptx(filepath):
    try:
        ppt = pptx.Presentation(filepath)
        text = []
        for slide in ppt.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return sanitize_text('\n'.join(text))
    except Exception as e:
        logging.warning(Fore.YELLOW + f'\n\n!! Error extracting text from PowerPoint file "{filepath}": {e} !!\n\n')
        return ""

def extract_text_from_pdf(filepath):
    try:
        text = []
        with pdfplumber.open(filepath) as pdf:
            for page in pdf.pages:
                page_text = page.extract_text()
                if page_text:
                    text.append(page_text)
        return sanitize_text('\n'.join(text))
    except Exception as e:
        logging.warning(Fore.YELLOW + f'\n\n!! Error extracting text from PDF file "{filepath}": {e} !!\n\n')
        return ""

def extract_text_from_docx(filepath):
    try:
        doc = DocxDocument(filepath)
        text = []
        for para in doc.paragraphs:
            text.append(para.text)
        return sanitize_text('\n'.join(text))
    except Exception as e:
        logging.warning(Fore.YELLOW + f'\n\n!! Error extracting text from DOCX file "{filepath}": {e} !!\n\n')
        return ""

def extract_text_from_txt(filepath):
    try:
        with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
            return sanitize_text(f.read())
    except Exception as e:
        logging.warning(Fore.YELLOW + f'\n\n!! Error extracting text from TXT file "{filepath}": {e} !!\n\n')
        return ""

def extract_text_from_doc(filepath):
    try:
        text = textract.process(filepath).decode('utf-8', errors='ignore')
        return sanitize_text(text)
    except Exception as e:
        logging.warning(Fore.YELLOW + f'\n\n!! Error extracting text from DOC file "{filepath}": {e} !!\n\n')
        return ""

def write_text_to_docx(text_dict, filename):
    doc = DocxDocument()
    for file_name, content in text_dict.items():
        doc.add_heading(file_name, level=1)
        doc.add_paragraph(content)
    doc.save(filename)

def find_files(directory, extensions):
    files_collected = []
    for root, dirs, files in os.walk(directory):
        for file in files:
            if file.endswith(tuple(extensions)):
                files_collected.append(os.path.join(root, file))
    return files_collected

# Specify the directory and file extensions to search
directory_to_search = '.'  # Start from the directory where the script is located
file_extensions = ['.pptx', '.pdf', '.docx', '.doc', '.txt']

# Find all files in the directory and its subdirectories with the specified extensions
files = find_files(directory_to_search, file_extensions)

# Add separation in logs
logging.info('\n' + '='*80 + '\n')
logging.info(f'Found {len(files)} files to process.')
logging.info('\n' + '='*80 + '\n')

# Initialize a dictionary to hold all extracted text
all_texts = {}

# Process each file based on its type with progress bar
for file in tqdm(files, desc="Processing files", unit="file"):
    if file.endswith('.pptx'):
        all_texts[os.path.basename(file)] = extract_text_from_pptx(file)
    elif file.endswith('.pdf'):
        all_texts[os.path.basename(file)] = extract_text_from_pdf(file)
    elif file.endswith('.docx'):
        all_texts[os.path.basename(file)] = extract_text_from_docx(file)
    elif file.endswith('.txt'):
        all_texts[os.path.basename(file)] = extract_text_from_txt(file)
    elif file.endswith('.doc'):
        all_texts[os.path.basename(file)] = extract_text_from_doc(file)

# Write all extracted text to a single Word document if there's text
if all_texts:
    write_text_to_docx(all_texts, 'combined_output.docx')
    logging.info('\n\n' + '='*80)
    logging.info(f'Finished processing. Extracted text saved to "combined_output.docx".')
    logging.info('='*80 + '\n\n')
else:
    logging.info('\n\n' + '='*80)
    logging.info('No documents were found or no text extracted, so no output file was created.')
    logging.info('='*80 + '\n\n')
